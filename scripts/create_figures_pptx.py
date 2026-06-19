"""
Create editable PPTX with all figures (one per slide, English).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pathlib import Path

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

figures = [
    {
        'file': 'figures/fig1_sharing_vs_distance.png',
        'title': 'Figure 1. Archaic DNA sharing vs. geographic distance',
        'caption': (
            'Scatter plot of pairwise archaic DNA segment sharing correlation vs. '
            'geographic distance for 66 populations (2,145 pairs). '
            '(A) Neanderthal: blue = same continent, red = cross-continent, '
            'orange triangles = admixed populations. Dashed lines show corrected '
            'regression. (B) Denisovan: purple diamonds = Oceania-involved pairs. '
            'Corrected for admixture and continental grouping (R² = 0.51 / 0.50).'
        ),
    },
    {
        'file': 'figures/fig2_sharing_heatmap.png',
        'title': 'Figure 2. Pairwise archaic DNA sharing heatmap',
        'caption': (
            'Heatmap of Neanderthal (left) and Denisovan (right) introgression '
            'segment sharing for 31 key populations. Note the isolated Oceanian '
            'cluster in the Denisovan panel (bottom-right), reflecting the Wallace '
            'Line boundary. Population labels colored by continental region.'
        ),
    },
    {
        'file': 'figures/fig3_minard_migration.png',
        'title': 'Figure 3. Minard-style Out-of-Africa migration flow',
        'caption': (
            'Flow visualization of Homo sapiens Out-of-Africa migration with archaic '
            'admixture events. Band width = relative population size. Stars = admixture '
            'events (Neanderthal ~47 kya; Denisovan 1 ~45 kya for Oceania 3-5%; '
            'Denisovan 2 ~30 kya for East Asia ~0.06%). Timelines approximate.'
        ),
    },
    {
        'file': 'figures/fig4_bivariate_world_map.png',
        'title': 'Figure 4. Bivariate world map of archaic DNA',
        'caption': (
            'Global distribution of archaic human DNA. Circle size = Neanderthal DNA '
            'proportion (0.08-1.8%). Color intensity = Denisovan DNA proportion '
            '(0.02-3.5%). Sharp color transition at the Wallace Line visible between '
            'mainland SE Asia and island Melanesia.'
        ),
    },
]

for fig_info in figures:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = fig_info['title']
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)

    # Image
    img_path = Path(fig_info['file'])
    if img_path.exists():
        # Scale to fit
        max_width = Inches(12.3)
        max_height = Inches(5.5)
        from PIL import Image
        with Image.open(str(img_path)) as img:
            w, h = img.size
        aspect = w / h
        if max_width / aspect <= max_height:
            img_w = max_width
            img_h = int(max_width / aspect)
        else:
            img_h = max_height
            img_w = int(max_height * aspect)

        left = (SLIDE_WIDTH - img_w) // 2
        top = Inches(0.9)
        slide.shapes.add_picture(str(img_path), left, top, width=img_w, height=img_h)
    else:
        print(f"Warning: {img_path} not found")

    # Caption
    cap_top = Inches(6.5)
    txBox2 = slide.shapes.add_textbox(Inches(0.5), cap_top, Inches(12.3), Inches(0.9))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = fig_info['caption']
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(80, 80, 80)

outpath = 'docs/figures_en.pptx'
prs.save(outpath)
print(f"PPTX saved to {outpath}")
